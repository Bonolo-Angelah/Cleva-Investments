"""
Cleva Investment - Professional PowerPoint Presentation Generator
Uses Gijima brand colors and creates a downloadable presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

# Gijima Brand Colors
GIJIMA_NAVY = RGBColor(30, 58, 138)      # #1e3a8a (Navy Blue)
GIJIMA_RED = RGBColor(220, 38, 38)       # #dc2626 (Red)
GIJIMA_LIGHT_BLUE = RGBColor(59, 130, 246)  # #3b82f6
GIJIMA_GRAY = RGBColor(107, 114, 128)    # #6b7280
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(243, 244, 246)

def add_footer(slide, logo_path=None):
    """Add footer to slide with Gijima logo"""
    # Add footer text box
    left = Inches(0.5)
    top = Inches(7)
    width = Inches(5)
    height = Inches(0.5)

    footer_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = footer_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = "Cleva Investment | AI-Powered Investment Advisory Platform"
    p.font.size = Pt(10)
    p.font.color.rgb = GIJIMA_GRAY

    # Add logo if provided
    if logo_path and os.path.exists(logo_path):
        try:
            slide.shapes.add_picture(logo_path, Inches(8.5), Inches(6.8), height=Inches(0.4))
        except:
            print(f"Warning: Could not add logo from {logo_path}")

def create_title_slide(prs, logo_path=None):
    """Create title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = GIJIMA_NAVY

    # Title
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(1.5)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = title_box.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = "CLEVA INVESTMENT"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(left, Inches(4), width, Inches(1))
    text_frame = subtitle_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = "AI-Powered Investment Advisory Platform"
    p.font.size = Pt(28)
    p.font.color.rgb = GIJIMA_LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Red accent line
    line_shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(3.5), Inches(5.2), Inches(3), Inches(0.05)
    )
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = GIJIMA_RED
    line_shape.line.fill.background()

    add_footer(slide, logo_path)

def create_section_slide(prs, title, logo_path=None):
    """Create section divider slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background gradient effect
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = GIJIMA_NAVY

    # Title
    left = Inches(1)
    top = Inches(3)
    width = Inches(8)
    height = Inches(1.5)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = title_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Accent bar
    accent_shape = slide.shapes.add_shape(
        1,
        Inches(0), Inches(4.5), Inches(10), Inches(0.1)
    )
    accent_shape.fill.solid()
    accent_shape.fill.fore_color.rgb = GIJIMA_RED
    accent_shape.line.fill.background()

    add_footer(slide, logo_path)

def create_content_slide(prs, title, content_items, logo_path=None):
    """Create content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title bar
    title_bar = slide.shapes.add_shape(
        1,
        Inches(0), Inches(0), Inches(10), Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = GIJIMA_NAVY
    title_bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    text_frame = title_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Content
    left = Inches(1)
    top = Inches(1.8)
    width = Inches(8)
    height = Inches(5)

    content_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, item in enumerate(content_items):
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]

        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = GIJIMA_GRAY
        p.level = 0
        p.space_before = Pt(12)

        # Add bullet point
        p.text = f"• {item}"

    add_footer(slide, logo_path)

def create_two_column_slide(prs, title, left_content, right_content, logo_path=None):
    """Create two-column content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title bar
    title_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = GIJIMA_NAVY
    title_bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    text_frame = title_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.5), Inches(5))
    text_frame = left_box.text_frame
    text_frame.word_wrap = True

    for i, item in enumerate(left_content):
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = GIJIMA_GRAY
        p.space_before = Pt(10)

    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.8), Inches(4), Inches(5))
    text_frame = right_box.text_frame
    text_frame.word_wrap = True

    for i, item in enumerate(right_content):
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = GIJIMA_GRAY
        p.space_before = Pt(10)

    add_footer(slide, logo_path)

def generate_presentation():
    """Generate the complete Cleva Investment presentation"""

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Path to logo (if exists) - looks in docs folder
    logo_path = os.path.join(os.path.dirname(__file__), 'gijima-logo.png')
    if not os.path.exists(logo_path):
        print("Note: Logo file not found at docs/gijima-logo.png")
        print("Please save the Gijima logo as 'gijima-logo.png' in the docs folder")
        logo_path = None

    # Slide 1: Title Slide
    create_title_slide(prs, logo_path)

    # Slide 2: Problem Statement Section
    create_section_slide(prs, "THE CHALLENGE", logo_path)

    # Slide 3: Problem Statement Details
    create_content_slide(prs, "Investment Challenges Facing South Africans", [
        "Lack of access to affordable, personalized financial advisory services",
        "Difficulty tracking and managing multiple investment portfolios",
        "Information overload making investment decisions overwhelming",
        "Challenges aligning investments with specific financial goals",
        "Limited tools designed for South African market and Rands (ZAR)",
        "High costs of traditional financial advisors exclude many investors"
    ], logo_path)

    # Slide 4: Solution Section
    create_section_slide(prs, "OUR SOLUTION", logo_path)

    # Slide 5: Solution Overview
    create_content_slide(prs, "Cleva Investment Platform", [
        "AI-powered investment advisory platform democratizing financial guidance",
        "Conversational AI chatbot providing personalized recommendations",
        "Comprehensive portfolio management and tracking tools",
        "Goal-based financial planning with progress monitoring",
        "Real-time market data integration for informed decisions",
        "Built specifically for South African investors with ZAR support"
    ], logo_path)

    # Slide 6: Key Features Section
    create_section_slide(prs, "KEY FEATURES", logo_path)

    # Slide 7: Core Features
    create_two_column_slide(prs, "Platform Capabilities", [
        "AI Investment Advisor (Cleva Bot)",
        "Natural language conversations",
        "Personalized recommendations",
        "Context-aware responses",
        "",
        "Portfolio Management",
        "Multiple portfolio support",
        "Real-time performance tracking",
        "Transaction management"
    ], [
        "Goal-Based Planning",
        "Custom financial goals",
        "Progress tracking",
        "Monthly contribution calculator",
        "",
        "Market Intelligence",
        "Live stock quotes",
        "Market news & articles",
        "Trending stocks & gainers"
    ], logo_path)

    # Slide 8: Analytics Features
    create_content_slide(prs, "Advanced Analytics & Insights", [
        "Sector Allocation Analysis - Visualize investment diversification",
        "Top Performers Dashboard - Identify best-performing assets",
        "Return on Investment (ROI) Metrics - Track portfolio performance",
        "Gain/Loss Tracking - Monitor profitability in real-time",
        "Multi-Currency Conversions - Automatic ZAR conversions for global assets",
        "Performance Cards - Quick overview of key metrics"
    ], logo_path)

    # Slide 9: Technology Section
    create_section_slide(prs, "TECHNOLOGY STACK", logo_path)

    # Slide 10: Technical Architecture
    create_two_column_slide(prs, "Robust & Scalable Architecture", [
        "Frontend",
        "React.js with modern hooks",
        "Tailwind CSS for responsive UI",
        "Socket.IO for real-time chat",
        "React Router for navigation",
        "",
        "Backend",
        "Node.js with Express.js",
        "RESTful API architecture",
        "Socket.IO for WebSockets"
    ], [
        "Databases",
        "PostgreSQL - User & portfolio data",
        "MongoDB - Chat history & articles",
        "Neo4j - Investment recommendations",
        "",
        "AI & Data",
        "Cohere/OpenAI integration",
        "Yahoo Finance API",
        "Financial Modeling Prep API"
    ], logo_path)

    # Slide 11: Value Proposition Section
    create_section_slide(prs, "VALUE PROPOSITION", logo_path)

    # Slide 12: Key Differentiators
    create_content_slide(prs, "Why Cleva Investment?", [
        "Democratized Access - Professional-grade tools for everyone",
        "AI-Powered Intelligence - Smart recommendations based on your profile",
        "Goal-Centric Approach - Every decision tied to your financial goals",
        "South African Focus - Built for ZAR and local market context",
        "Real-Time Data - Live market information for accurate insights",
        "All-in-One Platform - Advisory, management, and planning combined"
    ], logo_path)

    # Slide 13: User Benefits
    create_content_slide(prs, "Benefits for Investors", [
        "Save Money - Free alternative to expensive financial advisors",
        "Save Time - Automated tracking and intelligent recommendations",
        "Reduce Risk - Diversification insights and risk-aligned suggestions",
        "Stay Informed - Real-time market data and news",
        "Track Progress - Monitor goals and portfolio performance",
        "Make Better Decisions - Data-driven insights at your fingertips"
    ], logo_path)

    # Slide 14: Use Cases Section
    create_section_slide(prs, "USE CASES", logo_path)

    # Slide 15: Target Users
    create_two_column_slide(prs, "Who Benefits from Cleva?", [
        "Individual Investors",
        "First-time investors",
        "Young professionals",
        "Self-directed traders",
        "",
        "Goal-Oriented Savers",
        "Retirement planning",
        "Home purchase savings",
        "Education funds"
    ], [
        "Financial Enthusiasts",
        "Portfolio diversifiers",
        "Market followers",
        "Tech-savvy investors",
        "",
        "Budget-Conscious Users",
        "Those seeking free tools",
        "DIY investors",
        "Cost-sensitive savers"
    ], logo_path)

    # Slide 16: Future Roadmap Section
    create_section_slide(prs, "ROADMAP", logo_path)

    # Slide 17: Future Enhancements
    create_content_slide(prs, "Planned Enhancements", [
        "Mobile Applications - iOS and Android native apps",
        "Tax Optimization - Tax-loss harvesting and reporting",
        "Social Features - Community insights and shared strategies",
        "Advanced Analytics - Predictive models and forecasting",
        "Broker Integration - Direct trading capabilities",
        "Educational Content - Investment courses and tutorials"
    ], logo_path)

    # Slide 18: Closing Section
    create_section_slide(prs, "THANK YOU", logo_path)

    # Slide 19: Contact/Questions
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title bar
    title_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = GIJIMA_NAVY
    title_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    text_frame = title_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = "Questions?"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Content
    content_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(2))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = "Cleva Investment Platform"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = GIJIMA_NAVY
    p.alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = ""
    p.space_before = Pt(20)

    p = text_frame.add_paragraph()
    p.text = "AI-Powered Investment Advisory"
    p.font.size = Pt(18)
    p.font.color.rgb = GIJIMA_RED
    p.alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = ""
    p.space_before = Pt(20)

    p = text_frame.add_paragraph()
    p.text = "Making Professional Investment Guidance Accessible to All"
    p.font.size = Pt(16)
    p.font.color.rgb = GIJIMA_GRAY
    p.alignment = PP_ALIGN.CENTER

    add_footer(slide, logo_path)

    # Save presentation
    output_path = os.path.join(os.path.dirname(__file__), 'Cleva_Investment_Presentation.pptx')
    prs.save(output_path)
    print(f"[OK] Presentation created successfully!")
    print(f"[OK] Saved to: {output_path}")
    return output_path

if __name__ == "__main__":
    try:
        print("Generating Cleva Investment Presentation...")
        print("Using Gijima brand colors...")
        output_file = generate_presentation()
        print(f"\n{'='*60}")
        print(f"SUCCESS! Your presentation is ready for download.")
        print(f"{'='*60}")
        print(f"\nFile location: {output_file}")
        print("\nYou can now open and present this PowerPoint file!")
    except Exception as e:
        print(f"Error generating presentation: {e}")
        print("\nMake sure you have python-pptx installed:")
        print("  pip install python-pptx")
